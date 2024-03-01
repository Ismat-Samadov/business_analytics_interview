from monthly_expenses_generator import generate_monthly_expenses
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime
import os
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF

# Create a directory to store images if it doesn't exist
if not os.path.exists("images"):
    os.makedirs("images")

# Generate monthly expenses dataframe
df = generate_monthly_expenses()

# Convert birth_date and transaction_date to datetime
df['birth_date'] = pd.to_datetime(df['birth_date'])
df['transaction_date'] = pd.to_datetime(df['transaction_date'])

# Visualizations

# 1. Total Expenses by Category
total_expenses_category = df.groupby('category')['amount'].sum().reset_index()
plt.figure(figsize=(10, 6))
sns.barplot(x='category', y='amount', data=total_expenses_category)
plt.title('Total Expenses by Category')
plt.xlabel('Category')
plt.ylabel('Total Expenses')
filename1 = "total_expenses_by_category.png"
plt.savefig(os.path.join("images", filename1))
plt.close()

# 2. Most Expensive Subcategories
most_expensive_subcategories = df.groupby('subcategory')['amount'].sum().nlargest(5).reset_index()
plt.figure(figsize=(10, 6))
sns.barplot(x='amount', y='subcategory', data=most_expensive_subcategories)
plt.title('Top 5 Most Expensive Subcategories')
plt.xlabel('Total Expenses')
plt.ylabel('Subcategory')
filename2 = "most_expensive_subcategories.png"
plt.savefig(os.path.join("images", filename2))
plt.close()

# 3. Calculate transaction frequency for each part of the month
plt.figure(figsize=(10, 6))
early_month = df[df['transaction_date'].dt.day <= 10]['transaction_date'].dt.month.value_counts().sort_index()
mid_month = df[(df['transaction_date'].dt.day > 10) & (df['transaction_date'].dt.day <= 20)]['transaction_date'].dt.month.value_counts().sort_index()
late_month = df[df['transaction_date'].dt.day > 20]['transaction_date'].dt.month.value_counts().sort_index()
width = 0.25
plt.bar(early_month.index - width, early_month, width, label='Early Month')
plt.bar(mid_month.index, mid_month, width, label='Mid Month')
plt.bar(late_month.index + width, late_month, width, label='Late Month')
plt.title('Transaction Frequency Over Time (Divided by Month Parts)')
plt.xlabel('Month')
plt.ylabel('Frequency')
plt.xticks(early_month.index)
plt.legend()
filename3 = "transaction_frequency_over_time.png"
plt.savefig(os.path.join("images", filename3))
plt.close()

# 4. Average Transaction Amount by Category
average_transaction_category = df.groupby('category')['amount'].mean().reset_index()
plt.figure(figsize=(10, 6))
sns.barplot(x='category', y='amount', data=average_transaction_category)
plt.title('Average Transaction Amount by Category')
plt.xlabel('Category')
plt.ylabel('Average Amount')
filename4 = "average_transaction_amount_by_category.png"
plt.savefig(os.path.join("images", filename4))
plt.close()

# 5. Distribution of Transaction Amounts
plt.figure(figsize=(10, 6))
sns.histplot(df['amount'], bins=20, kde=True)
plt.title('Distribution of Transaction Amounts')
plt.xlabel('Amount')
plt.ylabel('Frequency')
filename5 = "distribution_of_transaction_amounts.png"
plt.savefig(os.path.join("images", filename5))
plt.close()

# 6. Customer Age Distribution
df['age'] = (datetime.now() - df['birth_date']).astype('<m8[Y]')
plt.figure(figsize=(10, 6))
sns.histplot(df['age'], bins=20, kde=True)
plt.title('Customer Age Distribution')
plt.xlabel('Age')
plt.ylabel('Frequency')
filename6 = "customer_age_distribution.png"
plt.savefig(os.path.join("images", filename6))
plt.close()

# 7. Most Visited Locations
most_visited_locations = df['location'].value_counts().reset_index()
plt.figure(figsize=(10, 6))
sns.barplot(x='index', y='location', data=most_visited_locations)
plt.title('Most Visited Locations')
plt.xlabel('Location')
plt.ylabel('Frequency')
filename7 = "most_visited_locations.png"
plt.savefig(os.path.join("images", filename7))
plt.close()

# 8. Transaction Amounts vs. Birth Year
plt.figure(figsize=(10, 6))
sns.scatterplot(x='birth_date', y='amount', data=df)
plt.title('Transaction Amounts vs. Birth Year')
plt.xlabel('Birth Year')
plt.ylabel('Amount')
filename8 = "transaction_amounts_vs_birth_year.png"
plt.savefig(os.path.join("images", filename8))
plt.close()

# 9. Transaction Amounts by Location
plt.figure(figsize=(10, 6))
sns.boxplot(x='location', y='amount', data=df)
plt.title('Transaction Amounts by Location')
plt.xlabel('Location')
plt.ylabel('Amount')
filename9 = "transaction_amounts_by_location.png"
plt.savefig(os.path.join("images", filename9))
plt.close()

# 10. Transaction Patterns
plt.figure(figsize=(10, 6))
sns.countplot(x='transaction_date', hue='category', data=df)
plt.title('Transaction Patterns')
plt.xlabel('Transaction Date')
plt.ylabel('Frequency')
filename10 = "transaction_patterns.png"
plt.savefig(os.path.join("images", filename10))
plt.close()

# Create PowerPoint presentation and insert images into slides
presentation = Presentation()

for image_file in sorted(os.listdir("images")):
    if image_file.endswith(".png"):
        slide_layout = presentation.slide_layouts[5]  # Choose appropriate slide layout
        slide = presentation.slides.add_slide(slide_layout)
        image_path = os.path.join("images", image_file)
        left = top = Inches(1)
        slide.shapes.add_picture(image_path, left, top, width=Inches(9), height=Inches(6))

presentation.save("monthly_expenses_analysis.pptx")

# Create PDF document and insert images
pdf = FPDF()
pdf.set_auto_page_break(auto=True, margin=15)
pdf.add_page()

for image_file in sorted(os.listdir("images")):
    if image_file.endswith(".png"):
        image_path = os.path.join("images", image_file)
        pdf.image(image_path, x=10, y=10, w=190)  # Assuming A4 page size, adjust coordinates as needed

pdf.output("monthly_expenses_analysis.pdf")

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import pandas as pd

df = pd.read_excel("cashback_analysis.xlsx")

plt.figure(figsize=(10, 6))

plt.subplot(1, 2, 1)
df_net_profit = df.drop(df.index[-1])  
plt.bar(df_net_profit["Cardholder ID"], df_net_profit["Net Profit/Loss"], color='blue')
plt.xlabel('Cardholder ID')
plt.ylabel('Net Profit/Loss ($)')
plt.title('Net Profit/Loss per Cardholder')
plt.savefig('net_profit_chart.png')  
plt.clf()  


plt.subplot(1, 2, 1)
cashback_labels = ['Supermarket Spending', 'Cashback Expenses']
cashback_sizes = [df['Supermarket Spending'].sum(), df['Total Cashback Expenses'].sum()]
plt.pie(cashback_sizes, labels=cashback_labels, autopct='%1.1f%%', startangle=140)
plt.axis('equal')
plt.title('Distribution of Cashback Expenses')
plt.savefig('cashback_expenses_chart.png')  
plt.clf()  

prs = Presentation()

slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
title.text = "Cashback Analysis"
subtitle = slide_1.placeholders[1]
subtitle.text = "Insights from Cashback Campaign"

slide_2 = prs.slides.add_slide(prs.slide_layouts[5])
title, content = slide_2.shapes.title, slide_2.placeholders[0]
title.text = "Net Profit/Loss per Cardholder"
content.text = "The bar plot illustrates the net profit/loss for each cardholder."
img_path_1 = 'net_profit_chart.png'
left, top, width, height = Inches(1), Inches(1.5), Inches(6), Inches(4.5)
slide_2.shapes.add_picture(img_path_1, left, top, width=width, height=height)

slide_3 = prs.slides.add_slide(prs.slide_layouts[5])
title, content = slide_3.shapes.title, slide_3.placeholders[0]
title.text = "Distribution of Cashback Expenses"
content.text = "The pie chart depicts the distribution of cashback expenses."
img_path_2 = 'cashback_expenses_chart.png'  
left, top, width, height = Inches(1), Inches(1.5), Inches(6), Inches(4.5)
slide_3.shapes.add_picture(img_path_2, left, top, width=width, height=height)

prs.save("cashback_analysis_presentation_with_insights.pptx")
print("Data analysis and insights have been saved to cashback_analysis_presentation_with_insights.pptx.")
